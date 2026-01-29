#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Extract development actions from PowerPoint file.

Extracts:
- Development actions per competency
- Actions grouped by level (1-5)
- Actions grouped by 70/20/10 logic
"""

import json
import re
import sys
from pathlib import Path
from pptx import Presentation

KNOWN_CLUSTERS = set()
KNOWN_BLOCKS = set()
KNOWN_COMPETENCIES = {}

def normalize_key(text):
    """Normalize text to create canonical keys for matching."""
    if not text:
        return ""
    text = str(text).strip()
    text = text.lower()
    text = re.sub(r'[^\w\s-]', '', text)
    text = re.sub(r'[-\s]+', '_', text)
    return text.strip('_')

EXCLUDED_ACTIONS = {
    normalize_key("Начать формулировать запросы на данные в структурированном виде: какая цель, какой вопрос нужно ответить, какие поля (колонки) нужны, за какой период, какие фильтры/разрезы, какой формат вывода."): True,
    normalize_key("Регулярно обсуждать с аналитиками, какие данные реально доступны, какие есть ограничения по качеству/частоте/точности, сколько времени займет реализация."): True,
    normalize_key("Принять участие в согласовании регулярного (еженедельного/ежемесячного) отчета: совместно определить, какие показатели и разрезы нужны для управления."): True,
    normalize_key("Проверять полученные выгрузки и отчеты: соответствуют ли они требованиям, нет ли ошибок."): True,
}

def clean_text(text, keep_linebreaks=True):
    if text is None:
        return ""
    value = str(text)
    value = value.replace("\x0b", "\n")
    value = value.replace("\r", "\n")
    if not keep_linebreaks:
        value = value.replace("\n", " ")
    value = re.sub(r"[ \t]+", " ", value)
    if keep_linebreaks:
        value = re.sub(r"\n+", "\n", value)
    return value.strip()

def extract_text_from_shape(shape):
    """Extract all text from a shape."""
    text_parts = []
    
    if hasattr(shape, "has_table") and shape.has_table:
        # Extract text from tables cell by cell
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text:
                    text_parts.append(cell.text)
    if hasattr(shape, "text"):
        text_parts.append(shape.text)
    elif hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    text_parts.append(run.text)
    
    return clean_text("\n".join(text_parts), keep_linebreaks=True)

def is_title_slide(text):
    """Check if slide is a title/contents slide."""
    title_keywords = [
        "содержание", "меню развивающих действий", "цель меню",
        "из чего состоит", "как пользоваться", "руководство",
        "словарь терминов", "перечень внешних образовательных ресурсов"
    ]
    text_lower = text.lower()
    return any(keyword in text_lower for keyword in title_keywords)

def is_resources_slide(text):
    """Check if slide contains resources list (not actions)."""
    text_lower = text.lower()
    return "список ресурсов" in text_lower or "перечень внешних образовательных ресурсов" in text_lower

def is_glossary_slide(text):
    """Check if slide contains glossary content."""
    return "словарь терминов" in text.lower()

def extract_competency_name_from_slide(slide):
    """Try to identify competency name from slide."""
    # Look for text that might be a competency name
    # Usually in title or first large text box
    
    if slide.shapes.title:
        title = clean_text(slide.shapes.title.text, keep_linebreaks=False)
        if title and not is_title_slide(title):
            title_key = normalize_key(title)
            if title_key in KNOWN_COMPETENCIES:
                return title
            if title_key in KNOWN_CLUSTERS or title_key in KNOWN_BLOCKS:
                return None
            if "обучение на практике" in title.lower() or "развитие на рабочем месте" in title.lower() or "обучение и саморазвитие" in title.lower():
                return None
            return title
    
    # Check first few text shapes
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            text = clean_text(extract_text_from_shape(shape), keep_linebreaks=False)
            if text and len(text) > 5 and len(text) < 100:
                text_shapes.append(text)
    
    # First substantial text might be competency name
    for text in text_shapes[:3]:
        text_lower = text.lower()
        if "обучение на практике" in text_lower or "развитие на рабочем месте" in text_lower or "обучение и саморазвитие" in text_lower:
            continue
        if re.search(r'\b70\s*%|\b20\s*%|\b10\s*%', text_lower):
            continue
        text_key = normalize_key(text)
        if text_key in KNOWN_COMPETENCIES:
            return KNOWN_COMPETENCIES.get(text_key, text)
        if text_key in KNOWN_CLUSTERS or text_key in KNOWN_BLOCKS:
            continue
        if not is_title_slide(text) and len(text) < 80:
            return text

    # Fallback: search all slide text for known competency names
    slide_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame or (hasattr(shape, "has_table") and shape.has_table):
            text = clean_text(extract_text_from_shape(shape), keep_linebreaks=False)
            if text:
                slide_texts.append(text)
    if slide_texts:
        joined_norm = normalize_key(" ".join(slide_texts))
        for comp_norm, comp_name in sorted(KNOWN_COMPETENCIES.items(), key=lambda x: len(x[0]), reverse=True):
            if comp_norm and comp_norm in joined_norm:
                return comp_name
    
    return None

def extract_actions_from_slide(slide, competency_name=None, initial_type=None, initial_level=None):
    """Extract development actions from a slide."""
    actions = []
    level = initial_level
    action_type = initial_type  # 70/20/10
    seen = set()
    
    all_text = []
    for shape in slide.shapes:
        if shape.has_text_frame or (hasattr(shape, "has_table") and shape.has_table):
            text = extract_text_from_shape(shape)
            if text:
                all_text.append((shape.top, shape.left, text))

    joined_text = "\n".join([t for _, _, t in all_text]).lower()
    has70 = bool(re.search(r'\b70\s*%|\b70\s*процентов|обучение на практике', joined_text))
    has20 = bool(re.search(r'\b20\s*%|\b20\s*процентов|развитие на рабочем месте', joined_text))
    has10 = bool(re.search(r'\b10\s*%|\b10\s*процентов|обучение и саморазвитие', joined_text))
    default_type = initial_type
    if default_type is None:
        if has70 and not has20 and not has10:
            default_type = "70"
        elif has20 and not has70 and not has10:
            default_type = "20"
        elif has10 and not has70 and not has20:
            default_type = "10"
    
    # Process all text to find actions (top-to-bottom, left-to-right)
    for _, _, text in sorted(all_text, key=lambda x: (x[0], x[1])):
        # Skip if this is the competency name
        if competency_name and normalize_key(text) == normalize_key(competency_name):
            continue

        # Extract action items
        lines = text.split('\n')
        for line in lines:
            raw_line = line.strip()
            is_bullet = bool(re.match(r'^[\u2022•\-\*]', raw_line) or re.match(r'^\d+[\.\)]', raw_line))
            line = clean_text(line, keep_linebreaks=False)

            line_lower = line.lower()

            if not line:
                continue

            # Update level if line indicates level
            level_match = re.search(r'(уровень|описание уровня)\s*(\d)', line_lower)
            if level_match:
                level = level_match.group(2)
                continue

            # Update action type if line indicates 70/20/10 section
            if re.search(r'\b70\s*%|\b70\s*процентов', line_lower) or "обучение на практике" in line_lower:
                action_type = "70"
                continue
            if re.search(r'\b20\s*%|\b20\s*процентов', line_lower) or "развитие на рабочем месте" in line_lower:
                action_type = "20"
                continue
            if re.search(r'\b10\s*%|\b10\s*процентов', line_lower) or "обучение и саморазвитие" in line_lower:
                action_type = "10"
                continue

            if len(line) < 10:
                continue
            
            # Remove bullet points and numbering
            line_clean = re.sub(r'^[\u2022•\-\*]+\s*', '', line)
            line_clean = re.sub(r'^\d+[\.\)]\s*', '', line_clean)
            line_clean = re.sub(r'^\d+\s*[-–]\s*', '', line_clean)
            
            # Skip if it's a header or very short
            if len(line_clean) < 10:
                continue
            
            # Skip common non-action text
            skip_patterns = [
                r'^уровень\s*\d',
                r'^описание\s*уровня',
                r'^целевой\s*уровень',
                r'^список\s*ресурсов',
                r'^ресурсы\s*для',
                r'^книга',
                r'^курс',
                r'^обучение\s+на\s+практике',
                r'^развитие\s+на\s+рабочем\s+месте',
                r'^обучение\s+и\s+саморазвитие',
                r'^\d+$'  # Just a number
            ]
            if any(re.match(pattern, line_clean.lower()) for pattern in skip_patterns):
                continue

            if "словарь терминов" in line_clean.lower() or "список ресурсов" in line_clean.lower():
                continue
            line_key = normalize_key(line_clean)
            if line_key in KNOWN_COMPETENCIES:
                continue
            if line_key in KNOWN_CLUSTERS or line_key in KNOWN_BLOCKS:
                continue
            if line_key in EXCLUDED_ACTIONS:
                continue

            if action_type is None:
                action_type = default_type or "70"

            # Merge short continuation lines into previous action when possible
            if actions and not is_bullet:
                prev_text = actions[-1].get("text", "")
                if prev_text and not re.search(r'[.!?]$', prev_text):
                    if len(line_clean.split()) <= 4:
                        actions[-1]["text"] = f"{prev_text} {line_clean}".strip()
                        continue

            if line_clean in seen:
                continue
            seen.add(line_clean)
            
            actions.append({
                "text": line_clean,
                "level": level,
                "type": action_type
            })
    
    return actions, action_type, level

def extract_all_actions(prs):
    """Extract all development actions from presentation."""
    all_actions = {}
    current_competency = None
    current_level = None
    current_type = None
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        # Skip title and resources slides
        slide_text = ""
        if slide.shapes.title:
            slide_text = clean_text(slide.shapes.title.text, keep_linebreaks=False)

        slide_all_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = extract_text_from_shape(shape)
                if text:
                    slide_all_text.append(text)

        joined_text = "\n".join(slide_all_text)

        skip = is_title_slide(slide_text) or is_resources_slide(joined_text) or is_glossary_slide(joined_text)
        if not skip and not slide_text:
            joined_lower = joined_text.lower()
            if "содержание" in joined_lower or "меню развивающих действий" in joined_lower:
                skip = True
        if skip:
            continue
        
        # Try to identify competency
        competency_name = extract_competency_name_from_slide(slide)
        
        if competency_name:
            # Check if this looks like a competency (not a block/cluster)
            # Competencies are usually shorter and more specific
            comp_key = normalize_key(competency_name)
            if comp_key in KNOWN_COMPETENCIES:
                competency_name = KNOWN_COMPETENCIES[comp_key]
            current_competency = competency_name
            current_level = None
            current_type = None
            if comp_key not in all_actions:
                all_actions[comp_key] = {
                    "competency_name": competency_name,
                    "actions_by_level": {},
                    "actions_by_type": {
                        "70": [],
                        "20": [],
                        "10": []
                    },
                    "all_actions": []
                }
        
        # Extract actions from this slide
        if current_competency:
            actions, current_type, current_level = extract_actions_from_slide(
                slide,
                current_competency,
                current_type,
                current_level
            )
            
            comp_key = normalize_key(current_competency)
            if comp_key in all_actions:
                for action in actions:
                    level = action.get("level") or "all"
                    action_type = action.get("type")
                    action_text = action.get("text")
                    
                    if level not in all_actions[comp_key]["actions_by_level"]:
                        all_actions[comp_key]["actions_by_level"][level] = []
                    
                    action_obj = {
                        "text": action_text,
                        "type": action_type
                    }
                    
                    all_actions[comp_key]["actions_by_level"][level].append(action_obj)
                    
                    if action_type and action_type in ["70", "20", "10"]:
                        all_actions[comp_key]["actions_by_type"][action_type].append(action_obj)
                    
                    all_actions[comp_key]["all_actions"].append(action_obj)
    
    return all_actions

def main():
    base_dir = Path(__file__).parent.parent
    pptx_path = base_dir / "data-src" / "menu.pptx"
    if not pptx_path.exists():
        pptx_path = base_dir / "data" / "Меню развивающих действий_ШЦР.pptx"
    
    if not pptx_path.exists():
        print(f"Error: PowerPoint file not found at {pptx_path}")
        sys.exit(1)
    
    # Load known cluster/block names to avoid mislabeling slides as competencies
    model_path = base_dir / "frontend" / "data" / "model.json"
    if model_path.exists():
        with open(model_path, "r", encoding="utf-8") as f:
            model_data = json.load(f)
        for c in model_data.get("clusters", []):
            KNOWN_CLUSTERS.add(normalize_key(c.get("name", "")))
        for b in model_data.get("blocks", []):
            KNOWN_BLOCKS.add(normalize_key(b.get("name", "")))
        for comp in model_data.get("competencies", []):
            comp_name = comp.get("name", "")
            comp_norm = normalize_key(comp_name)
            if comp_norm:
                KNOWN_COMPETENCIES[comp_norm] = comp_name

    print(f"Loading PowerPoint file: {pptx_path}")
    prs = Presentation(pptx_path)
    
    print(f"Total slides: {len(prs.slides)}")
    
    # Extract actions
    all_actions = extract_all_actions(prs)
    
    # Save output
    output_path = Path(__file__).parent.parent / "frontend" / "data" / "actions.json"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(all_actions, f, ensure_ascii=False, indent=2)
    
    print(f"\n✓ Extracted actions saved to: {output_path}")
    print(f"  - Competencies with actions: {len(all_actions)}")
    
    # Print summary
    for comp_key, comp_data in list(all_actions.items())[:10]:
        print(f"\n  {comp_data['competency_name']}:")
        print(f"    - Total actions: {len(comp_data['all_actions'])}")
        levels = [k for k in comp_data['actions_by_level'].keys() if k != 'all']
        if levels:
            print(f"    - Levels: {levels}")
        type_counts = {k: len(v) for k, v in comp_data['actions_by_type'].items() if v}
        if type_counts:
            print(f"    - 70/20/10: {type_counts}")
    
    return all_actions

if __name__ == "__main__":
    main()
