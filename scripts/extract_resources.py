#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Extract resources lists from PowerPoint and map them to clusters.
Outputs frontend/data/resources.json
"""

import json
import re
from pathlib import Path
from pptx import Presentation

def normalize(text):
    if not text:
        return ""
    return re.sub(r"[^\w]+", "", text.lower())

def extract_texts(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            texts.append(shape.text.strip())
    return texts

def is_resources_slide(text):
    return "список ресурсов" in text.lower()

def extract_resources_for_slide(texts, cluster_title):
    lines = []
    for t in texts:
        for line in t.splitlines():
            line = line.strip()
            if not line:
                continue
            if "список ресурсов" in line.lower():
                continue
            if normalize(line) == normalize(cluster_title):
                continue
            if re.match(r"^\d+[\.\)]", line) or "—" in line:
                line = re.sub(r"^\d+[\.\)]\s*", "", line)
                lines.append(line)

    seen = set()
    result = []
    for line in lines:
        if line in seen:
            continue
        seen.add(line)
        match = re.search(r"https?://\S+", line)
        result.append({
            "text": line,
            "url": match.group(0) if match else ""
        })
    return result

def main():
    project_dir = Path(__file__).parent.parent
    model_path = project_dir / "frontend" / "data" / "model.json"
    pptx_path = project_dir / "data" / "Меню развивающих действий_ШЦР.pptx"
    output_path = project_dir / "frontend" / "data" / "resources.json"

    if not model_path.exists():
        raise SystemExit(f"Model file not found: {model_path}")
    if not pptx_path.exists():
        raise SystemExit(f"PPTX file not found: {pptx_path}")

    model = json.load(open(model_path, "r", encoding="utf-8"))
    clusters = model.get("clusters", [])
    cluster_norm = {c["id"]: normalize(c["name"]) for c in clusters}

    prs = Presentation(str(pptx_path))
    resources_by_cluster = {}

    for slide in prs.slides:
        texts = extract_texts(slide)
        if not texts:
            continue
        joined = "\n".join(texts)
        if not is_resources_slide(joined):
            continue

        candidates = [t.splitlines()[0].strip() for t in texts if t.strip()]
        cluster_id = None
        cluster_title = None
        for cand in candidates:
            cand_norm = normalize(cand)
            for cid, cnorm in cluster_norm.items():
                if cand_norm and (cand_norm in cnorm or cnorm in cand_norm):
                    cluster_id = cid
                    cluster_title = cand
                    break
            if cluster_id:
                break

        if not cluster_id:
            continue

        resources_by_cluster[cluster_id] = extract_resources_for_slide(texts, cluster_title)

    output_path.write_text(
        json.dumps({"resources_by_cluster": resources_by_cluster}, ensure_ascii=False, indent=2),
        encoding="utf-8"
    )

    missing = [c["id"] for c in clusters if c["id"] not in resources_by_cluster]
    print(f"Resources extracted for clusters: {len(resources_by_cluster)}")
    if missing:
        print(f"Warning: no resources for clusters: {missing}")

if __name__ == "__main__":
    main()
